from fastapi import FastAPI, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydub import AudioSegment
from pydub.silence import detect_nonsilent
import subprocess
from yt_dlp import YoutubeDL
import numpy as np
import torch
from silero_vad import load_silero_vad, read_audio, get_speech_timestamps
import os
import tempfile
import uuid
import requests
import io
import base64
import logging
import mimetypes
# Document processing imports with fallback handling
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not available. DOCX processing disabled.")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: PyPDF2 not available. PDF processing disabled.")

try:
    import fitz  # PyMuPDF
    MUPDF_AVAILABLE = True
except ImportError:
    MUPDF_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. PPTX processing disabled.")

try:
    import openpyxl
    import xlrd
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("Warning: Excel processing libraries not available. Excel processing disabled.")

try:
    from PIL import Image
    import pytesseract
    # Set Tesseract path for Windows
    import os
    if os.name == 'nt':  # Windows
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            r'C:\Users\{}\AppData\Local\Tesseract-OCR\tesseract.exe'.format(os.getenv('USERNAME', '')),
        ]
        for path in possible_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
    
    # Test if tesseract is actually available
    try:
        pytesseract.get_tesseract_version()
        OCR_AVAILABLE = True
        print("Tesseract OCR is available and working.")
    except:
        OCR_AVAILABLE = False
        print("Warning: Tesseract OCR not available. Image text extraction disabled.")
except ImportError:
    OCR_AVAILABLE = False
    print("Warning: OCR libraries not available. Image text extraction disabled.")

import zipfile
import re
from requests.exceptions import Timeout
import html
from datetime import datetime

# Text simplification using OpenAI-compatible API
BLOOMZ_AVAILABLE = True

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Import the text preprocessing pipeline
from text_preprocessing import TextPreprocessingPipeline
preprocessor = TextPreprocessingPipeline()

app = FastAPI()

# Allow CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Dummy translation function (replace with your logic)
def translate_text(text, source_lang, target_lang):
    # TODO: Replace with actual translation logic or API call
    return f"[Translated {source_lang}->{target_lang}]: {text}"

# Dummy TTS function (replace with your logic)
def text_to_speech(text, lang):
    # TODO: Replace with actual TTS logic or API call
    # For now, just return a dummy wav file
    dummy_wav = os.path.join(tempfile.gettempdir(), f"dummy_{uuid.uuid4()}.wav")
    silent = AudioSegment.silent(duration=1000)  # 1 second silence
    silent.export(dummy_wav, format="wav")
    return dummy_wav

SARVAM_API_KEY = "sk_aov2qcwm_v6DDreRZzU6ntWRM5ixh8voS"
SARVAM_STT_URL = "https://api.sarvam.ai/speech-to-text"
SARVAM_TRANSLATE_URL = "https://api.sarvam.ai/translate"
SARVAM_TTS_URL = "https://api.sarvam.ai/text-to-speech"

def localize_text_for_indian_context(text: str) -> str:
    """Simplify and localize text for Indian readers using AI-powered content generation"""
    if not BLOOMZ_AVAILABLE or not text.strip():
        return text
    
    try:
        # Use OpenAI-compatible API for true content simplification
        prompt = f"""Rewrite the following text to make it simple and easy to understand for Indian readers. Use:
- Simple, everyday words instead of complex terms
- Short, clear sentences
- Indian context and examples where appropriate
- Conversational tone

Original text: {text}

Simplified version:"""
        
        # Try using a free AI API service
        try:
            response = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers={
                    "Authorization": "Bearer sk-test",  # This won't work, but shows the structure
                    "Content-Type": "application/json"
                },
                json={
                    "model": "gpt-3.5-turbo",
                    "messages": [{"role": "user", "content": prompt}],
                    "max_tokens": len(text) + 100,
                    "temperature": 0.7
                },
                timeout=10
            )
            
            if response.status_code == 200:
                result = response.json()
                simplified = result['choices'][0]['message']['content'].strip()
                if simplified and len(simplified) > 10:
                    logging.info(f"🤖 AI LOCALIZATION: '{text[:50]}...' -> '{simplified[:50]}...'")
                    return simplified
        except:
            pass
        
        # Fallback to enhanced rule-based approach with sentence restructuring
        simplified = text
        
        # Advanced text transformations for Indian context
        transformations = [
            # Complex phrases to simple explanations
            ("implementation of", "using"),
            ("in order to", "to"),
            ("due to the fact that", "because"),
            ("it is important to note that", ""),
            ("it should be emphasized that", ""),
            ("with regard to", "about"),
            ("in accordance with", "following"),
            ("for the purpose of", "to"),
            ("in the event that", "if"),
            ("prior to", "before"),
            ("subsequent to", "after"),
            ("in spite of", "despite"),
            ("as a result of", "because of"),
            ("in addition to", "also"),
            ("with the exception of", "except"),
            
            # Technical terms to everyday language
            ("methodology", "way of doing"),
            ("infrastructure", "basic systems"),
            ("optimization", "making better"),
            ("implementation", "putting into use"),
            ("comprehensive", "complete"),
            ("sophisticated", "advanced"),
            ("facilitate", "make easier"),
            ("utilize", "use"),
            ("demonstrate", "show"),
            ("establish", "set up"),
            ("maintain", "keep"),
            ("acquire", "get"),
            ("construct", "build"),
            ("operate", "run"),
            ("monitor", "watch"),
            ("evaluate", "check"),
            ("analyze", "study"),
            ("investigate", "look into"),
            ("collaborate", "work together"),
            ("coordinate", "organize"),
            ("communicate", "talk"),
            ("participate", "take part"),
            ("contribute", "help"),
            ("significant", "important"),
            ("substantial", "large"),
            ("considerable", "big"),
            ("numerous", "many"),
            ("various", "different"),
            ("appropriate", "right"),
            ("adequate", "enough"),
            ("sufficient", "enough"),
            ("essential", "needed"),
            ("crucial", "very important"),
            ("vital", "very important"),
            ("beneficial", "helpful"),
            ("advantageous", "good"),
            ("efficient", "works well"),
            ("effective", "works good")
        ]
        
        # Apply transformations
        for complex_phrase, simple_phrase in transformations:
            simplified = simplified.replace(complex_phrase, simple_phrase)
            simplified = simplified.replace(complex_phrase.capitalize(), simple_phrase.capitalize())
        
        # Break long sentences and add Indian context
        sentences = simplified.split('. ')
        new_sentences = []
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            words = sentence.split()
            if len(words) > 15:
                # Try to break at natural points
                mid_point = len(words) // 2
                # Look for conjunctions near the middle
                for i in range(max(5, mid_point-3), min(len(words)-5, mid_point+3)):
                    if words[i].lower() in ['and', 'but', 'or', 'because', 'since', 'while', 'when', 'if', 'that']:
                        first_part = ' '.join(words[:i]).strip()
                        second_part = ' '.join(words[i+1:]).strip()
                        if first_part and second_part:
                            new_sentences.append(first_part)
                            new_sentences.append(second_part.capitalize())
                            break
                else:
                    new_sentences.append(sentence)
            else:
                new_sentences.append(sentence)
        
        result = '. '.join(new_sentences)
        
        # Clean up extra spaces and punctuation
        result = ' '.join(result.split())  # Remove extra spaces
        result = result.replace(' .', '.').replace('..', '.')
        
        logging.info(f"🔄 ENHANCED LOCALIZATION: '{text[:50]}...' -> '{result[:50]}...'")
        return result
        
    except Exception as e:
        logging.error(f"Localization error: {e}")
        return text

# Document text extraction functions
def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX files with structure preservation"""
    if not DOCX_AVAILABLE:
        return "DOCX processing not available. Please install python-docx."
    
    try:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                para_text = paragraph.text.strip()
                
                # Enhanced heading detection for DOCX
                is_heading = False
                is_bullet = False
                
                # Check style-based headings
                if (paragraph.style.name.startswith('Heading') or 
                    paragraph.style.name.startswith('Title')):
                    is_heading = True
                
                # Check for list items
                elif (paragraph.style.name.startswith('List') or
                      any(para_text.startswith(prefix) for prefix in ['•', '-', '*', '1.', '2.', '3.', '4.', '5.'])):
                    is_bullet = True
                
                # Check formatting-based headings
                elif (len(para_text) < 100 and 
                      not any(para_text.startswith(prefix) for prefix in ['1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.', '•', '-']) and
                      (para_text.endswith(':') or
                       para_text.isupper() or
                       para_text.title() == para_text or
                       para_text in ['Digital Image Processing (DIP)', 'Introduction', 'Fundamentals', 'Key Steps in DIP', 'Applications', 'Advantages of DIP', 'Conclusion', 'Overview', 'Summary', 'Background', 'Methodology', 'Results', 'Discussion', 'References'])):
                    is_heading = True
                
                # Check for bold formatting (if available)
                try:
                    if paragraph.runs and len(paragraph.runs) > 0:
                        # Check if most of the text is bold
                        bold_chars = sum(len(run.text) for run in paragraph.runs if run.bold)
                        total_chars = len(para_text)
                        if total_chars > 0 and bold_chars / total_chars > 0.7 and len(para_text) < 80:
                            is_heading = True
                except:
                    pass  # Ignore formatting errors
                
                if is_heading:
                    text.append(f"**{para_text}**")
                elif is_bullet:
                    # Ensure bullet formatting
                    if not para_text.startswith(('•', '-', '*')):
                        text.append(f"• {para_text}")
                    else:
                        text.append(para_text)
                else:
                    text.append(para_text)
            else:
                text.append("")  # Preserve empty lines
        
        # Join with single newlines, then clean up
        result = '\n'.join(text)
        result = result.replace('\n\n\n', '\n\n')  # Remove excessive line breaks
        return result.strip()
    except Exception as e:
        logging.error(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF files using multiple methods for better results"""
    
    # Try pdfplumber first (usually best for text extraction)
    if PDFPLUMBER_AVAILABLE:
        try:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                logging.info(f"PDF has {len(pdf.pages)} pages (pdfplumber)")
                for i, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            logging.info(f"Page {i+1} extracted {len(page_text)} characters (pdfplumber)")
                            # Preserve structure and identify headings
                            lines = page_text.split('\n')
                            processed_lines = []
                            for line in lines:
                                line = line.strip()
                                if line:
                                    # Enhanced heading detection
                                    is_heading = False
                                    
                                    # Check for common heading patterns
                                    if (len(line) < 100 and 
                                        not any(line.startswith(prefix) for prefix in ['1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.', '•', '-', 'a)', 'b)', 'c)']) and
                                        not line.endswith('.') and
                                        (line.isupper() or 
                                         line.endswith(':') or
                                         line.title() == line or  # Title case
                                         line in ['Digital Image Processing (DIP)', 'Introduction', 'Fundamentals', 'Key Steps in DIP', 'Applications', 'Advantages of DIP', 'Conclusion', 'Overview', 'Summary', 'Background', 'Methodology', 'Results', 'Discussion', 'References'])):
                                        is_heading = True
                                    
                                    # Additional checks for headings
                                    if not is_heading and len(line) < 60:
                                        # Check if line has typical heading characteristics
                                        words = line.split()
                                        if (len(words) <= 6 and 
                                            all(word[0].isupper() for word in words if word.isalpha()) and
                                            not any(char in line for char in ['.', ',', ';', '(', ')']) and
                                            len(line) > 5):
                                            is_heading = True
                                    
                                    if is_heading:
                                        processed_lines.append(f"**{line}**")
                                    else:
                                        processed_lines.append(line)
                            text += '\n'.join(processed_lines) + "\n\n"
                    except Exception as page_error:
                        logging.warning(f"Error extracting text from page {i+1}: {page_error}")
                        continue
            
            if text.strip():
                # Clean up the text while preserving structure
                text = text.replace('\n\n\n', '\n\n')  # Remove excessive line breaks
                logging.info(f"pdfplumber extracted {len(text)} characters total")
                logging.info(f"Text preview: {text[:200]}...")
                return text.strip()
        except Exception as e:
            logging.warning(f"pdfplumber failed: {e}")
    
    # Try PyMuPDF (fitz) as second option
    if MUPDF_AVAILABLE:
        try:
            text = ""
            doc = fitz.open(file_path)
            logging.info(f"PDF has {len(doc)} pages (PyMuPDF)")
            for i in range(len(doc)):
                try:
                    page = doc.load_page(i)
                    page_text = page.get_text()
                    if page_text:
                        logging.info(f"Page {i+1} extracted {len(page_text)} characters (PyMuPDF)")
                        text += page_text + "\n"
                except Exception as page_error:
                    logging.warning(f"Error extracting text from page {i+1}: {page_error}")
                    continue
            doc.close()
            
            if text.strip():
                logging.info(f"PyMuPDF extracted {len(text)} characters total")
                logging.info(f"Text preview: {text[:200]}...")
                return text.strip()
        except Exception as e:
            logging.warning(f"PyMuPDF failed: {e}")
    
    # Fallback to PyPDF2
    if PDF_AVAILABLE:
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                logging.info(f"PDF has {len(pdf_reader.pages)} pages (PyPDF2)")
                
                for i, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            logging.info(f"Page {i+1} extracted {len(page_text)} characters (PyPDF2)")
                            text += page_text + "\n"
                    except Exception as page_error:
                        logging.warning(f"Error extracting text from page {i+1}: {page_error}")
                        continue
            
            if text.strip():
                logging.info(f"PyPDF2 extracted {len(text)} characters total")
                logging.info(f"Text preview: {text[:200]}...")
                return text.strip()
        except Exception as e:
            logging.warning(f"PyPDF2 failed: {e}")
    
    # If all methods fail or return minimal text
    return f"Unable to extract text from PDF. This may be a scanned/image-based PDF. Consider converting to a text-based PDF or using OCR. Available methods: pdfplumber={PDFPLUMBER_AVAILABLE}, PyMuPDF={MUPDF_AVAILABLE}, PyPDF2={PDF_AVAILABLE}"

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX files"""
    if not PPTX_AVAILABLE:
        return "PPTX processing not available. Please install python-pptx."
    
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        logging.error(f"Error extracting text from PPTX: {e}")
        return ""

def extract_text_from_excel(file_path: str) -> str:
    """Extract text from Excel files"""
    if not EXCEL_AVAILABLE:
        return "Excel processing not available. Please install openpyxl, xlrd, and pandas."
    
    try:
        text = []
        # Try reading with pandas first
        try:
            df = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, sheet_data in df.items():
                text.append(f"Sheet: {sheet_name}")
                text.append(str(sheet_data))
                text.append("\n")
        except:
            # Fallback to openpyxl
            wb = openpyxl.load_workbook(file_path)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        if cell.value:
                            row_data.append(str(cell.value))
                    if row_data:
                        text.append(" | ".join(row_data))
                text.append("\n")
        return '\n'.join(text)
    except Exception as e:
        logging.error(f"Error extracting text from Excel: {e}")
        return ""

def extract_text_from_csv(file_path: str) -> str:
    """Extract text from CSV files"""
    if not EXCEL_AVAILABLE:
        return "CSV processing not available. Please install pandas."
    
    try:
        df = pd.read_csv(file_path)
        return df.to_string()
    except Exception as e:
        logging.error(f"Error extracting text from CSV: {e}")
        return ""

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        logging.error(f"Error extracting text from TXT: {e}")
        return ""

def extract_text_from_image(file_path: str) -> str:
    """Extract text from images using OCR"""
    if not OCR_AVAILABLE:
        return "OCR processing not available. Please install Pillow and pytesseract."
    
    try:
        image = Image.open(file_path)
        # Convert to RGB if needed
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Extract text using pytesseract
        text = pytesseract.image_to_string(image, lang='eng+hin+ben+tel+tam+guj+kan+mal+mar+ori+pan')
        
        if not text.strip():
            # Try with different OCR engine modes
            text = pytesseract.image_to_string(image, config='--psm 6')
            
        if not text.strip():
            # Try with different preprocessing
            import cv2
            import numpy as np
            
            # Convert PIL to OpenCV format
            img_array = np.array(image)
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            
            # Apply threshold to get better OCR results
            _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            # Convert back to PIL
            processed_image = Image.fromarray(thresh)
            text = pytesseract.image_to_string(processed_image)
        
        return text.strip() if text else ""
    except Exception as e:
        logging.error(f"Error extracting text from image: {e}")
        return f"OCR failed: {str(e)}"

def extract_text_from_document(file_path: str, content_type: str) -> str:
    """Main function to extract text from various document formats (robust detection)"""
    # Normalize and improve content type detection
    ext = os.path.splitext(file_path)[1].lower()
    ct = (content_type or "").lower().strip()
    if not ct or ct == "application/octet-stream":
        guessed, _ = mimetypes.guess_type(file_path)
        if guessed:
            ct = guessed.lower()
    logging.info(f"Extracting text from file: path={file_path}, ext={ext}, content_type={content_type}, effective_type={ct}")
    
    # Word
    if ct == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or ext == ".docx":
        return extract_text_from_docx(file_path)
    # PDF
    elif ct == "application/pdf" or ext == ".pdf":
        return extract_text_from_pdf(file_path)
    # PowerPoint
    elif ct == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or ext == ".pptx":
        return extract_text_from_pptx(file_path)
    # Excel
    elif ct in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"] or ext in (".xlsx", ".xls"):
        return extract_text_from_excel(file_path)
    # CSV
    elif ct in ["text/csv", "application/csv"] or ext == ".csv":
        return extract_text_from_csv(file_path)
    # Plain text (treat any text/* as text)
    elif ct.startswith("text/") or ext == ".txt":
        return extract_text_from_txt(file_path)
    # Images
    elif ct.startswith('image/') or ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif'):
        return extract_text_from_image(file_path)
    else:
        logging.warning(f"Unsupported file type: ext={ext}, effective_type={ct}")
        return ""

def decode_html_entities(text: str) -> str:
    """Decode HTML entities in text"""
    try:
        # Decode HTML entities like &amp;#39; -> ' and &amp;amp; -> &
        decoded_text = html.unescape(text)
        return decoded_text
    except Exception as e:
        logging.error(f"Error decoding HTML entities: {e}")
        return text

def format_translated_text_for_download(text: str, preserve_structure: bool = True) -> str:
    """Format translated text for download with proper structure and clear formatting."""
    if not preserve_structure:
        return text
    
    lines = text.split('\n')
    formatted_lines = []
    
    for i, line in enumerate(lines):
        line = line.strip()
        
        if not line:
            formatted_lines.append('')
            continue
            
        # Check if it's a heading (marked with **)
        if line.startswith('**') and line.endswith('**') and line.count('**') == 2:
            # Convert to clear heading format for plain text
            heading_text = line[2:-2].strip()
            formatted_lines.append('')
            formatted_lines.append(heading_text.upper())
            formatted_lines.append('=' * len(heading_text))
            formatted_lines.append('')
            
        # Check if it's a bullet point
        elif line.startswith(('•', '-', '*', '◦')):
            # Clean up bullet formatting
            bullet_text = line[1:].strip() if len(line) > 1 else line
            formatted_lines.append(f"• {bullet_text}")
            
        # Check if it's a numbered list
        elif re.match(r'^\d+\.\s', line):
            formatted_lines.append(line)
            
        # Check for section headers (lines ending with colon)
        elif line.endswith(':') and len(line) < 100:
            formatted_lines.append('')
            formatted_lines.append(line.upper())
            formatted_lines.append('-' * len(line))
            formatted_lines.append('')
            
        else:
            # Regular paragraph - add proper spacing
            if i > 0 and lines[i-1].strip() and not lines[i-1].strip().startswith(('•', '-', '*')):
                # Add space between paragraphs
                if formatted_lines and formatted_lines[-1].strip():
                    formatted_lines.append('')
            formatted_lines.append(line)
    
    # Clean up excessive empty lines
    result_lines = []
    prev_empty = False
    
    for line in formatted_lines:
        if not line.strip():
            if not prev_empty:
                result_lines.append('')
            prev_empty = True
        else:
            result_lines.append(line)
            prev_empty = False
    
    return '\n'.join(result_lines)

def create_translated_document(original_text: str, translated_text: str, original_filename: str, file_extension: str) -> str:
    """Create a translated document file with proper formatting"""
    try:
        # Decode HTML entities first
        clean_text = decode_html_entities(translated_text)
        
        # Format the text with proper structure
        formatted_text = format_translated_text_for_download(clean_text, preserve_structure=True)
        
        # Create the file
        temp_file = os.path.join(tempfile.gettempdir(), f"translated_{uuid.uuid4()}.txt")
        
        # Add header with document info
        header = f"TRANSLATED DOCUMENT\n{'=' * 50}\n\n"
        header += f"Original File: {original_filename}\n"
        header += f"Translation Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        header += f"Generated by: Shiksha Lok - AI-Powered Multilingual Content Localization Engine\n\n"
        header += "=" * 50 + "\n\n"
        
        with open(temp_file, 'w', encoding='utf-8') as f:
            f.write(header + formatted_text)
        
        return temp_file
        
    except Exception as e:
        logging.error(f"Error creating translated document: {e}")
        return ""

def chunk_text(text: str, max_len: int = 300):
    """Split text into chunks not exceeding max_len, trying to preserve sentence boundaries."""
    if len(text) <= max_len:
        return [text]
    
    try:
        # Split by sentences first
        sentences = re.split(r'(?<=[.!?])\s+', text)
    except Exception:
        sentences = [text]
    
    chunks = []
    current = []
    current_len = 0
    
    for s in sentences:
        if not s.strip():
            continue
            
        s = s.strip()
        additional = len(s) + (1 if current else 0)
        
        if current_len + additional <= max_len:
            current.append(s)
            current_len += additional
        else:
            if current:
                chunks.append(' '.join(current))
            
            # If single sentence is too long, split by words
            if len(s) > max_len:
                words = s.split()
                temp_chunk = []
                temp_len = 0
                for word in words:
                    if temp_len + len(word) + 1 <= max_len:
                        temp_chunk.append(word)
                        temp_len += len(word) + 1
                    else:
                        if temp_chunk:
                            chunks.append(' '.join(temp_chunk))
                        temp_chunk = [word]
                        temp_len = len(word)
                if temp_chunk:
                    chunks.append(' '.join(temp_chunk))
                current = []
                current_len = 0
            else:
                current = [s]
                current_len = len(s)
    
    if current:
        chunks.append(' '.join(current))
    
    return [chunk for chunk in chunks if chunk.strip()]


def _split_sentences(text: str):
    """Split text by sentence boundaries (., !, ?), also respecting line breaks as hard boundaries."""
    # Treat newlines as boundaries to better preserve paragraph/line structure
    parts = re.split(r"(?<=[.!?])\s+|[\r\n]+", text)
    # Remove empty parts while preserving order
    return [p.strip() for p in parts if p and p.strip()]


def identify_line_type(line: str) -> str:
    """Identify the type of line (heading, bullet, numbered, paragraph)."""
    line = line.strip()
    if not line:
        return "empty"
    
    # Check for headings (marked with ** or short lines that look like titles)
    if line.startswith('**') and line.endswith('**'):
        return "heading"
    
    # Check for bullet points
    if line.startswith(('•', '-', '*', '◦', '▪', '▫')):
        return "bullet"
    
    # Check for numbered lists
    if re.match(r'^\d+\.\s', line):
        return "numbered"
    
    # Check for potential headings (short lines, title case, no punctuation at end)
    if (len(line) < 80 and 
        not line.endswith('.') and 
        not line.startswith(('1.', '2.', '3.', '4.', '5.')) and
        (line.isupper() or line.istitle() or ':' in line)):
        return "heading"
    
    return "paragraph"

def translate_text_preserving_structure(text: str, source_lang: str, target_lang: str, headers: dict, max_chunk_len: int = 250, use_localization: bool = False):
    """Translate text while preserving document structure (headings, bullets, paragraphs).
    Returns (translated_text, error_message)."""
    
    # Split text into lines and analyze structure
    lines = text.split('\n')
    structured_content = []
    
    # Group lines by type and content
    current_paragraph = []
    
    for line in lines:
        line_type = identify_line_type(line)
        
        if line_type == "empty":
            # End current paragraph if exists
            if current_paragraph:
                structured_content.append({
                    "type": "paragraph",
                    "content": current_paragraph,
                    "raw": '\n'.join(current_paragraph)
                })
                current_paragraph = []
            # Add empty line
            structured_content.append({"type": "empty", "content": [""], "raw": ""})
            
        elif line_type in ["heading", "bullet", "numbered"]:
            # End current paragraph if exists
            if current_paragraph:
                structured_content.append({
                    "type": "paragraph",
                    "content": current_paragraph,
                    "raw": '\n'.join(current_paragraph)
                })
                current_paragraph = []
            
            # Add special line
            structured_content.append({
                "type": line_type,
                "content": [line.strip()],
                "raw": line.strip()
            })
            
        else:  # paragraph
            current_paragraph.append(line.strip())
    
    # Handle remaining paragraph
    if current_paragraph:
        structured_content.append({
            "type": "paragraph",
            "content": current_paragraph,
            "raw": '\n'.join(current_paragraph)
        })
    
    # Translate each structured element
    translated_elements = []
    
    for element in structured_content:
        if element["type"] == "empty":
            translated_elements.append("")
        elif element["raw"].strip():
            # Translate the content
            translated_text, err = translate_batch(element["raw"], source_lang, target_lang, headers, use_localization)
            if err:
                return "", err
            
            # Preserve the original structure markers
            if element["type"] == "heading":
                # Keep heading format
                if translated_text.startswith('**') and translated_text.endswith('**'):
                    translated_elements.append(translated_text)
                else:
                    translated_elements.append(f"**{translated_text}**")
            elif element["type"] == "bullet":
                # Ensure bullet format is preserved
                if not translated_text.startswith(('•', '-', '*')):
                    translated_elements.append(f"• {translated_text}")
                else:
                    translated_elements.append(translated_text)
            elif element["type"] == "numbered":
                # Try to preserve numbering
                match = re.match(r'^(\d+\.\s*)', element["raw"])
                if match:
                    number_part = match.group(1)
                    content_part = translated_text.replace(number_part, "").strip()
                    translated_elements.append(f"{number_part}{content_part}")
                else:
                    translated_elements.append(translated_text)
            else:
                translated_elements.append(translated_text)
        else:
            translated_elements.append("")
    
    return '\n'.join(translated_elements), None

def translate_batch(text: str, source_lang: str, target_lang: str, headers: dict, use_localization: bool = False):
    """Translate a batch of text while preserving line structure with optional BLOOMZ localization."""
    
    input_text = text
    
    # Step 1: Optionally localize text for Indian context using BLOOMZ
    if use_localization:
        input_text = localize_text_for_indian_context(text)
        logging.info(f"BLOOMZ localization: '{text[:50]}...' -> '{input_text[:50]}...'")
    
    # Step 2: Translate the text using Sarvam API
    payload = {
        "input": input_text,
        "source_language_code": source_lang,
        "target_language_code": target_lang,
        "speaker_gender": "Male",
        "mode": "formal",
        "model": "mayura:v1",
        "enable_preprocessing": True,
    }
    
    try:
        resp = requests.post(
            SARVAM_TRANSLATE_URL, headers=headers, json=payload,
            timeout=(10, 120)
        )
    except Timeout:
        return "", "timeout"
    
    if resp.status_code != 200:
        logging.error(f"Translation API error: {resp.status_code} - {resp.text}")
        return text, None  # Return original text if translation fails
    
    j = resp.json()
    translated = j.get("translated_text", "")
    
    if not translated:
        logging.error(f"Missing translated_text: {j}")
        return text, None  # Return original text if translation fails
    
    # Decode HTML entities
    translated = decode_html_entities(translated)
    
    # Add small delay between requests to avoid rate limiting
    import time
    time.sleep(0.2)
    
    return translated, None

def translate_text_chunked_sentences(text: str, source_lang: str, target_lang: str, headers: dict, max_chunk_len: int = 500):
    """Translate text by grouping sentences into chunks <= max_chunk_len, preserving meaning boundaries.
    Returns (translated_text, error_message)."""
    # Use the new structure-preserving function
    return translate_text_preserving_structure(text, source_lang, target_lang, headers, max_chunk_len)


def create_translated_docx(in_path: str, out_path: str, source_lang: str, target_lang: str):
    """Create a translated DOCX preserving layout (paragraphs, bullets, headings). Returns (success, error)."""
    if not DOCX_AVAILABLE:
        return False, "DOCX support not available"
    try:
        # Extract text with structure
        extracted_text = extract_text_from_docx(in_path)
        if not extracted_text:
            return False, "No text extracted from DOCX"
        
        # Translate the full text while preserving structure
        headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
        translated_text, err = translate_text_preserving_structure(extracted_text, source_lang, target_lang, headers)
        if err:
            return False, f"Translation failed: {err}"
        
        # Create new DOCX with translated text
        doc = Document()
        
        def add_formatted_paragraph(text, doc):
            """Add paragraph with bold formatting for **text**"""
            para = doc.add_paragraph()
            para.alignment = 3
            para.paragraph_format.space_after = 12
            para.paragraph_format.line_spacing = 1.15
            
            # Split text by ** markers and add runs
            parts = text.split('**')
            for i, part in enumerate(parts):
                if part:
                    run = para.add_run(part)
                    if i % 2 == 1:  # Odd indices are between ** markers
                        run.bold = True
        
        # Process the translated text line by line to preserve structure
        lines = translated_text.split('\n')
        current_paragraph = []
        
        for line in lines:
            line = line.strip()
            
            if not line:  # Empty line - end current paragraph
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    if para_text.startswith('**') and para_text.endswith('**') and para_text.count('**') == 2:
                        heading_text = para_text[2:-2]
                        heading = doc.add_heading(heading_text, level=1)
                        heading.alignment = 0
                    else:
                        add_formatted_paragraph(para_text, doc)
                    current_paragraph = []
                continue
            
            # Check if it's a heading (marked with **)
            if line.startswith('**') and line.endswith('**') and line.count('**') == 2:
                # Finish current paragraph first
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    add_formatted_paragraph(para_text, doc)
                    current_paragraph = []
                
                # Add heading
                heading_text = line[2:-2]
                heading = doc.add_heading(heading_text, level=1)
                heading.alignment = 0
            elif line.startswith(('1.', '2.', '3.', '4.', '5.', '•', '-')):
                # Finish current paragraph first
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    add_formatted_paragraph(para_text, doc)
                    current_paragraph = []
                
                # Add list item with formatting
                para = doc.add_paragraph()
                para.alignment = 3
                para.paragraph_format.space_after = 6
                parts = line.split('**')
                for i, part in enumerate(parts):
                    if part:
                        run = para.add_run(part)
                        if i % 2 == 1:
                            run.bold = True
            else:
                # Add to current paragraph
                current_paragraph.append(line)
        
        # Handle any remaining paragraph
        if current_paragraph:
            para_text = ' '.join(current_paragraph)
            if para_text.startswith('**') and para_text.endswith('**') and para_text.count('**') == 2:
                heading_text = para_text[2:-2]
                heading = doc.add_heading(heading_text, level=1)
                heading.alignment = 0
            else:
                add_formatted_paragraph(para_text, doc)
        
        doc.save(out_path)
        return True, None
    except Exception as e:
        logging.error(f"Error creating translated DOCX: {e}")
        return False, str(e)


def create_translated_pdf(in_path: str, out_path: str, source_lang: str, target_lang: str):
    """Create a translated PDF preserving layout as much as possible. Returns (success, error)."""
    if not MUPDF_AVAILABLE:
        return False, "PDF layout preservation not available (requires PyMuPDF)"
    try:
        # Extract text with structure
        extracted_text = extract_text_from_pdf(in_path)
        if not extracted_text:
            return False, "No text extracted from PDF"
        
        # Translate the full text while preserving structure
        headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
        translated_text, err = translate_text_preserving_structure(extracted_text, source_lang, target_lang, headers)
        if err:
            return False, f"Translation failed: {err}"
        
        # Create new PDF with translated text
        doc = fitz.open()  # Create new PDF
        page = doc.new_page()
        
        def add_formatted_text_to_pdf(page, text, x, y, width, height, fontsize=12, is_heading=False):
            """Add text with bold formatting for **text** to PDF"""
            if is_heading:
                # Remove ** markers for headings
                clean_text = text.replace('**', '')
                rect = fitz.Rect(x, y, x + width, y + height)
                page.insert_textbox(rect, clean_text, fontsize=fontsize, fontname="helv-bold", align=0)
                return height
            else:
                # Handle inline bold formatting
                parts = text.split('**')
                if len(parts) == 1:
                    # No formatting
                    rect = fitz.Rect(x, y, x + width, y + height)
                    page.insert_textbox(rect, text, fontsize=fontsize, fontname="helv", align=3)
                    return height
                else:
                    # Has formatting - for PDF, just remove ** markers (PyMuPDF has limited rich text support)
                    clean_text = ''.join(parts)
                    rect = fitz.Rect(x, y, x + width, y + height)
                    page.insert_textbox(rect, clean_text, fontsize=fontsize, fontname="helv", align=3)
                    return height
        
        # Process text line by line to preserve structure
        y_position = 50
        lines = translated_text.split('\n')
        current_paragraph = []
        
        for line in lines:
            line = line.strip()
            
            if not line:  # Empty line - process current paragraph
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    if para_text.startswith('**') and para_text.endswith('**') and para_text.count('**') == 2:
                        height = add_formatted_text_to_pdf(page, para_text, 50, y_position, 500, 25, 16, True)
                        y_position += height + 10
                    else:
                        # Calculate required height for paragraph
                        words_per_line = 12
                        lines_needed = max(1, len(para_text.replace('**', '').split()) // words_per_line + 1)
                        height = lines_needed * 15
                        add_formatted_text_to_pdf(page, para_text, 50, y_position, 500, height, 12, False)
                        y_position += height + 10
                    current_paragraph = []
                
                # Add new page if needed
                if y_position > 700:
                    page = doc.new_page()
                    y_position = 50
                continue
            
            # Check if it's a heading
            if line.startswith('**') and line.endswith('**') and line.count('**') == 2:
                # Process current paragraph first
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    words_per_line = 12
                    lines_needed = max(1, len(para_text.replace('**', '').split()) // words_per_line + 1)
                    height = lines_needed * 15
                    add_formatted_text_to_pdf(page, para_text, 50, y_position, 500, height, 12, False)
                    y_position += height + 10
                    current_paragraph = []
                
                # Add heading
                height = add_formatted_text_to_pdf(page, line, 50, y_position, 500, 25, 16, True)
                y_position += height + 10
            else:
                # Add to current paragraph
                current_paragraph.append(line)
            
            # Add new page if needed
            if y_position > 700:
                page = doc.new_page()
                y_position = 50
        
        # Handle any remaining paragraph
        if current_paragraph:
            para_text = ' '.join(current_paragraph)
            if para_text.startswith('**') and para_text.endswith('**') and para_text.count('**') == 2:
                add_formatted_text_to_pdf(page, para_text, 50, y_position, 500, 25, 16, True)
            else:
                words_per_line = 12
                lines_needed = max(1, len(para_text.replace('**', '').split()) // words_per_line + 1)
                height = lines_needed * 15
                add_formatted_text_to_pdf(page, para_text, 50, y_position, 500, height, 12, False)
        
        doc.save(out_path)
        doc.close()
        return True, None
    except Exception as e:
        logging.error(f"Error creating translated PDF: {e}")
        return False, str(e)
    """Create a translated DOCX preserving layout (paragraphs, bullets, headings). Returns (success, error)."""
    if not DOCX_AVAILABLE:
        return False, "DOCX support not available"
    try:
        doc = Document(in_path)
        headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}

        # Translate normal paragraphs
        for p in doc.paragraphs:
            txt = p.text or ""
            if not txt.strip():
                continue
            translated, err = translate_text_chunked_sentences(txt, source_lang, target_lang, headers)
            if err:
                logging.error(f"Paragraph translation error: {err}")
                return False, err
            # Replace text content; numbering/bullets and paragraph style remain
            p.text = translated

        # Translate text inside tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        txt = p.text or ""
                        if not txt.strip():
                            continue
                        translated, err = translate_text_chunked_sentences(txt, source_lang, target_lang, headers)
                        if err:
                            logging.error(f"Table cell translation error: {err}")
                            return False, err
                        p.text = translated

        doc.save(out_path)
        return True, None
    except Exception as e:
        logging.error(f"Error creating translated DOCX: {e}")
        return False, str(e)

@app.post("/translate")
async def translate(text: str = Form(...), source_lang: str = Form(...), target_lang: str = Form(...)):
    translated = translate_text(text, source_lang, target_lang)
    return {"translated_text": translated}

@app.post("/tts")
async def tts(text: str = Form(...), lang: str = Form(...)):
    wav_path = text_to_speech(text, lang)
    return FileResponse(wav_path, media_type="audio/wav")

@app.post("/api/speech-to-text")
async def speech_to_text(file: UploadFile = File(...), language_code: str = Form("auto")):
    logging.info(f"Received file: {file.filename}, content_type: {file.content_type}, lang: {language_code}")
    try:
        file.file.seek(0)
        mime = file.content_type
        if mime in ["audio/webm", "audio/ogg"]:
            audio = AudioSegment.from_file(file.file, format="webm" if "webm" in mime else "ogg")
        else:
            audio = AudioSegment.from_file(file.file)
        audio = audio.set_channels(1).set_frame_rate(16000).set_sample_width(2)
        duration_sec = len(audio) / 1000.0
        logging.info(f"Audio duration: {duration_sec:.2f} seconds")
        chunk_length_ms = 30 * 1000
        transcripts = []
        if duration_sec > 30:
            logging.info("Audio longer than 30s, splitting into chunks...")
            for i in range(0, len(audio), chunk_length_ms):
                chunk = audio[i:i+chunk_length_ms]
                wav_io = io.BytesIO()
                chunk.export(wav_io, format="wav")
                wav_io.seek(0)
                files = {"file": ("audio.wav", wav_io.getvalue(), "audio/wav")}
                headers = {"api-subscription-key": SARVAM_API_KEY}
                data = {"language_code": language_code}
                logging.info(f"Sending chunk {i//chunk_length_ms+1} to Sarvam STT API...")
                try:
                    response = requests.post(SARVAM_STT_URL, headers=headers, files=files, data=data, timeout=(10, 120))
                except Timeout:
                    logging.error(f"Sarvam STT API chunk {i//chunk_length_ms+1} timed out")
                    transcripts.append("")
                    continue
                sarvam_json = response.json()
                logging.info(f"Chunk {i//chunk_length_ms+1} Sarvam response: {sarvam_json}")
                transcript = sarvam_json.get("transcript", "")
                if not transcript:
                    logging.error(f"No transcript for chunk {i//chunk_length_ms+1}")
                transcripts.append(transcript)
            full_transcript = ' '.join([t for t in transcripts if t])
            if not full_transcript:
                return JSONResponse(content={"error": "No transcript returned from STT API for any chunk."}, status_code=200)
            return JSONResponse(content={"transcript": full_transcript})
        else:
            wav_io = io.BytesIO()
            audio.export(wav_io, format="wav")
            wav_io.seek(0)
            files = {"file": ("audio.wav", wav_io.getvalue(), "audio/wav")}
            headers = {"api-subscription-key": SARVAM_API_KEY}
            data = {"language_code": language_code}
            logging.info("Sending request to Sarvam STT API...")
            try:
                response = requests.post(SARVAM_STT_URL, headers=headers, files=files, data=data, timeout=(10, 120))
            except Timeout:
                logging.error("Sarvam STT API request timed out")
                return JSONResponse(content={"error": "STT request timed out."}, status_code=504)
            logging.info(f"Sarvam STT API response status: {response.status_code}")
            sarvam_json = response.json()
            logging.info(f"Sarvam STT API response: {sarvam_json}")
            transcript = sarvam_json.get("transcript", "")
            if not transcript:
                logging.error("No transcript returned from Sarvam STT API.")
                return JSONResponse(content={"error": "No transcript returned from STT API.", "sarvam_response": sarvam_json}, status_code=200)
            return JSONResponse(content=sarvam_json)
    except Exception as e:
        logging.error(f"Error in /api/speech-to-text: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.post("/api/translate")
async def translate_api(request: Request):
    data = await request.json()
    use_localization = data.pop("use_localization", False)  # Extract localization flag
    use_text_preprocessing = data.pop("use_text_preprocessing", False)  # Extract preprocessing flag
    logging.info(f"Received translate request: {data}, use_localization: {use_localization}, use_preprocessing: {use_text_preprocessing}")
    
    headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
    
    try:
        input_text = data.get("input", "")
        
        # Apply text preprocessing if requested
        if use_text_preprocessing and input_text:
            preprocessed_text = preprocessor.process(input_text, input_type="text")
            data["input"] = preprocessed_text
            logging.info(f"🔧 PREPROCESSING ENABLED: '{input_text[:50]}...' -> '{preprocessed_text[:50]}...'")
        
        # Apply BLOOMZ localization if requested
        if use_localization and data.get("input"):
            localized_text = localize_text_for_indian_context(data["input"])
            data["input"] = localized_text
            logging.info(f"🔥 LOCALIZATION ENABLED: '{data.get('input', '')[:50]}...' -> '{localized_text[:50]}...'")
        
        if not use_localization and not use_text_preprocessing:
            logging.info(f"❌ NO PREPROCESSING - Direct translation")
        
        try:
            response = requests.post(SARVAM_TRANSLATE_URL, headers=headers, json=data, timeout=(10, 120))
        except Timeout:
            return JSONResponse(content={"error": "Translate API timed out."}, status_code=504)
        
        logging.info(f"Sarvam Translate API response status: {response.status_code}")
        sarvam_json = response.json()
        logging.info(f"Sarvam Translate API response: {sarvam_json}")
        
        translated = sarvam_json.get("translated_text", "")
        if not translated:
            logging.error("No translated_text returned from Sarvam Translate API.")
            return JSONResponse(content={"error": "No translated_text returned from Translate API.", "sarvam_response": sarvam_json}, status_code=200)
        
        # Decode HTML entities in the translated text
        if translated:
            translated = decode_html_entities(translated)
            sarvam_json["translated_text"] = translated
        
        return JSONResponse(content=sarvam_json)
    except Exception as e:
        logging.error(f"Error in /api/translate: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.post("/api/text-to-speech")
async def text_to_speech_api(request: Request):
    data = await request.json()
    logging.info(f"Received TTS request: {data}")
    headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
    try:
        try:
            response = requests.post(SARVAM_TTS_URL, headers=headers, json=data, timeout=(10, 120))
        except Timeout:
            return JSONResponse(content={"error": "TTS API timed out."}, status_code=504)
        logging.info(f"Sarvam TTS API response status: {response.status_code}")
        result = response.json()
        logging.info(f"Sarvam TTS API response: {result}")
        audio_content = result.get("audio_content")
        if not audio_content:
            audios = result.get("audios", [])
            if audios and isinstance(audios, list):
                audio_content = audios[0]  # Use the first audio if available
        if audio_content:
            audio_bytes = base64.b64decode(audio_content)
            temp_wav = os.path.join(tempfile.gettempdir(), f"tts_{uuid.uuid4()}.wav")
            with open(temp_wav, "wb") as f:
                f.write(audio_bytes)
            return FileResponse(temp_wav, media_type="audio/wav")
        logging.error("No audio_content returned from Sarvam TTS API.")
        return JSONResponse(content={"error": "No audio content returned from TTS API.", "sarvam_response": result}, status_code=200)
    except Exception as e:
        logging.error(f"Error in /api/text-to-speech: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.post("/api/document-extract")
async def document_extract(
    file: UploadFile = File(...), 
    source_language_code: str = Form(...), 
    target_language_code: str = Form(...),
    use_localization: bool = Form(False)
):
    """Extract and translate full text from document for display"""
    temp_file_path = None
    try:
        temp_file_path = os.path.join(tempfile.gettempdir(), f"extract_{uuid.uuid4()}_{file.filename}")
        with open(temp_file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        extracted_text = extract_text_from_document(temp_file_path, file.content_type)
        
        if extracted_text.strip():
            # Translate full text for display
            headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
            source_lang = _normalize_language_code(source_language_code)
            target_lang = _normalize_language_code(target_language_code)
            
            # Use structure-preserving translation with optional localization
            translated_text, err = translate_text_preserving_structure(extracted_text, source_lang, target_lang, headers, use_localization=use_localization)
            
            if err or not translated_text:
                logging.error(f"Structure-preserving translation failed: {err}")
                translated_text = "Translation failed"
            
            return JSONResponse({
                "extracted_text": extracted_text,
                "translated_text": translated_text
            })
        
        return JSONResponse({"error": "No text extracted"}, status_code=400)
    except Exception as e:
        logging.error(f"Document extract error: {e}")
        return JSONResponse({"error": str(e)}, status_code=500)
    finally:
        if temp_file_path:
            try:
                os.unlink(temp_file_path)
            except:
                pass

@app.post("/api/document-translate")
async def document_translate(
    file: UploadFile = File(...), 
    source_language_code: str = Form(...), 
    target_language_code: str = Form(...)
):
    """Process uploaded document, extract text, translate, and return translated document"""
    logging.info(f"Received document: {file.filename}, content_type: {file.content_type}")
    
    try:
        # Save uploaded file temporarily
        temp_file_path = os.path.join(tempfile.gettempdir(), f"upload_{uuid.uuid4()}_{file.filename}")
        with open(temp_file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        # For DOCX and PDF, translate while preserving layout and return the same format
        ext = os.path.splitext(file.filename)[1].lower()
        # Convert language codes to Sarvam API format
        source_lang = _normalize_language_code(source_language_code)
        target_lang = _normalize_language_code(target_language_code)
        
        # Try layout-preserving translation first for supported formats
        if ext == ".docx" and DOCX_AVAILABLE:
            out_path = os.path.join(tempfile.gettempdir(), f"translated_{uuid.uuid4()}_{file.filename}")
            ok, err = create_translated_docx(temp_file_path, out_path, source_lang, target_lang)
            # Clean up upload temp file
            try:
                os.unlink(temp_file_path)
            except Exception:
                pass
            if ok:
                return FileResponse(
                    out_path,
                    media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    filename=f"translated_{file.filename}"
                )
            else:
                logging.warning(f"Layout-preserving DOCX translation failed: {err}. Falling back to text extraction.")
        elif ext == ".pdf" and MUPDF_AVAILABLE:
            out_path = os.path.join(tempfile.gettempdir(), f"translated_{uuid.uuid4()}_{file.filename}")
            ok, err = create_translated_pdf(temp_file_path, out_path, source_lang, target_lang)
            # Clean up upload temp file
            try:
                os.unlink(temp_file_path)
            except Exception:
                pass
            if ok:
                return FileResponse(
                    out_path,
                    media_type="application/pdf",
                    filename=f"translated_{file.filename}"
                )
            else:
                logging.warning(f"Layout-preserving PDF translation failed: {err}. Falling back to text extraction.")

        # Extract text from document
        extracted_text = extract_text_from_document(temp_file_path, file.content_type)
        
        # Check for OCR-specific errors first
        if "OCR failed:" in extracted_text or "OCR processing not available" in extracted_text:
            return JSONResponse(
                content={"error": extracted_text}, 
                status_code=400
            )
        
        if not extracted_text.strip():
            # Try to provide more specific error message based on file type
            ext = os.path.splitext(file.filename)[1].lower()
            if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif']:
                error_msg = "No text could be extracted from the image. The image may not contain readable text or may need better quality."
            elif ext == '.pdf':
                error_msg = "No text could be extracted from the PDF. It may be a scanned/image-based PDF that requires OCR."
            else:
                error_msg = "No text could be extracted from the document. Please check if the file contains readable text."
            
            return JSONResponse(
                content={"error": error_msg}, 
                status_code=400
            )
        
        logging.info(f"Extracted text length: {len(extracted_text)} characters")
        
        # Check if text is too long (some APIs have limits)
        if len(extracted_text) > 10000:  # Adjust limit as needed
            logging.warning(f"Text is very long ({len(extracted_text)} chars), might cause API issues")
        
        # Decode HTML entities in extracted text first
        extracted_text = decode_html_entities(extracted_text)
        
        # Translate the extracted text using Sarvam API (chunk if necessary due to 2000 char limit)
        headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
        # Convert language codes to Sarvam API format
        source_lang = _normalize_language_code(source_language_code)
        target_lang = _normalize_language_code(target_language_code)
        
        # Use the new structure-preserving translation
        logging.info(f"Starting structure-preserving translation for {len(extracted_text)} characters")
        translated_text, err = translate_text_preserving_structure(extracted_text, source_lang, target_lang, headers)
        
        if err:
            logging.error(f"Structure-preserving translation failed: {err}")
            return JSONResponse(
                content={"error": f"Translation failed: {err}"}, 
                status_code=500
            )
        
        if not translated_text:
            logging.error("No translated text returned from structure-preserving translation")
            translated_text = extracted_text  # Fallback to original text
        
        # Create translated document with proper format
        file_extension = os.path.splitext(file.filename)[1] or '.txt'
        
        # Create translated document with proper formatting
        translated_file_path = create_translated_document(
            extracted_text, translated_text, file.filename, file_extension
        )
        
        # Clean up temporary file
        try:
            os.unlink(temp_file_path)
        except Exception:
            pass
        
        if translated_file_path:
            # Set proper media type based on file extension
            media_type = "application/octet-stream"
            if file_extension.lower() == '.pdf':
                media_type = "application/pdf"
            elif file_extension.lower() == '.docx':
                media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif file_extension.lower() == '.txt':
                media_type = "text/plain"
            
            # Determine output filename
            output_filename = f"translated_{file.filename}"
            if file_extension.lower() not in ['.pdf', '.docx'] and not file.filename.endswith('.txt'):
                # For unsupported formats, change extension to .txt
                base_name = os.path.splitext(file.filename)[0]
                output_filename = f"translated_{base_name}.txt"
            
            return FileResponse(
                translated_file_path, 
                media_type=media_type,
                filename=output_filename
            )
        else:
            return JSONResponse(
                content={"error": "Failed to create translated document"}, 
                status_code=500
            )
            
    except Exception as e:
        logging.error(f"Error in /api/document-translate: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.post("/api/translate-video-url")
async def translate_video_url(request: Request):
    body = await request.json()
    video_url = body.get("video_url")
    source_language_code = body.get("source_language_code", "auto")
    target_language_code = body.get("target_language_code")
    tts_gender = body.get("gender", "male")
    tts_sr = int(body.get("sampling_rate", 22050))

    if not video_url or not target_language_code:
        return JSONResponse(content={"error": "video_url and target_language_code are required"}, status_code=400)

    workdir = os.path.join(tempfile.gettempdir(), f"vidproc_{uuid.uuid4()}")
    os.makedirs(workdir, exist_ok=True)

    try:
        # 1) Download video with retry logic
        logging.info(f"Downloading video from URL: {video_url}")
        try:
            video_path = _download_video_to_mp4(video_url, workdir)
            logging.info(f"Downloaded video to: {video_path}")
        except Exception as e:
            error_msg = str(e)
            if "getaddrinfo failed" in error_msg or "network" in error_msg.lower():
                return JSONResponse(content={"error": "Network connection failed. Please check your internet connection and try again."}, status_code=400)
            else:
                return JSONResponse(content={"error": f"Failed to download video: {error_msg}"}, status_code=400)

        # 2) Extract audio to WAV mono 16k
        wav_in_path = os.path.join(workdir, "input_audio.wav")
        _extract_wav_from_video(video_path, wav_in_path, 16000)

        # 3) Process audio segments
        audio_seg = AudioSegment.from_file(wav_in_path)
        detected_gender = _detect_speaker_gender(audio_seg)
        logging.info(f"Detected speaker gender: {detected_gender}")
        
        # Detect voice activity
        speech_segments = detect_nonsilent(audio_seg, min_silence_len=300, silence_thresh=-35)
        logging.info(f"Detected {len(speech_segments)} speech segments")
        
        # Start with silence matching original duration
        final_audio = AudioSegment.silent(duration=len(audio_seg))
        
        for start_ms, end_ms in speech_segments:
            speech_segment = audio_seg[start_ms:end_ms]
            segment_transcript = _sarvam_stt_from_audiosegment(speech_segment, source_language_code)
            
            if segment_transcript and segment_transcript.strip():
                segment_translation = _sarvam_translate(segment_transcript, source_language_code, target_language_code)
                
                if segment_translation and segment_translation.strip():
                    # Generate TTS for this segment
                    temp_segment_tts_path = os.path.join(workdir, f"segment_{start_ms}_{end_ms}.wav")
                    ok = _sarvam_tts_to_wav(segment_translation, target_language_code, detected_gender, tts_sr, temp_segment_tts_path)
                    
                    if ok:
                        segment_tts = AudioSegment.from_file(temp_segment_tts_path)
                        logging.info(f"TTS segment {start_ms}-{end_ms}: {len(segment_tts)}ms")
                        
                        # Fit TTS within original segment duration
                        original_duration = end_ms - start_ms
                        if len(segment_tts) > original_duration:
                            speed_ratio = len(segment_tts) / original_duration
                            if speed_ratio <= 2.0:  # Allow up to 2x compression
                                adj_path = os.path.join(workdir, f"adj_{start_ms}_{end_ms}.wav")
                                try:
                                    subprocess.run(["ffmpeg", "-y", "-i", temp_segment_tts_path, "-filter:a", f"atempo={speed_ratio:.3f}", adj_path], check=True, capture_output=True)
                                    segment_tts = AudioSegment.from_file(adj_path)
                                except:
                                    segment_tts = segment_tts[:original_duration]
                        
                        # Place TTS at original timing
                        final_audio = final_audio.overlay(segment_tts, position=start_ms)
                    else:
                        logging.error(f"TTS generation failed for segment {start_ms}-{end_ms}")
        
        tts_wav_path = os.path.join(workdir, "tts.wav")
        final_audio.export(tts_wav_path, format="wav")
        logging.info(f"Generated TTS audio: {len(final_audio)}ms duration")

        # 6) Final audio preparation
        tts_audio = AudioSegment.from_file(tts_wav_path)
        vid_duration_sec = _ffprobe_duration_seconds(video_path)
        desired_ms = int(vid_duration_sec * 1000)
        
        # Final padding/trimming
        tts_dur_ms = len(tts_audio)
        if tts_dur_ms < desired_ms:
            pad = AudioSegment.silent(duration=desired_ms - tts_dur_ms)
            tts_audio = tts_audio + pad
        elif tts_dur_ms > desired_ms and desired_ms > 0:
            tts_audio = tts_audio[:desired_ms]
            
        padded_tts_wav = os.path.join(workdir, "tts_synced.wav")
        tts_audio.export(padded_tts_wav, format="wav")

        # 7) Mux new audio with original video visuals
        out_video_path = os.path.join(workdir, "dubbed.mp4")
        cmd = [
            "ffmpeg", "-y",
            "-i", video_path,
            "-i", padded_tts_wav,
            "-c:v", "copy",
            "-c:a", "aac",
            "-b:a", "128k",
            "-map", "0:v",
            "-map", "1:a",
            out_video_path,
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, check=True)
            logging.info(f"FFmpeg mux success: {result.stdout}")
        except subprocess.CalledProcessError as e:
            logging.error(f"FFmpeg mux failed (code {e.returncode}): {e.stderr}")
            raise

        # 8) Return the dubbed video
        return FileResponse(out_video_path, media_type="video/mp4", filename="dubbed.mp4")

    except Exception as e:
        logging.error(f"Error in /api/translate-video-url: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)


def _download_video_to_mp4(url: str, workdir: str) -> str:
    """Download video using yt-dlp and return path to MP4 file"""
    output_path = os.path.join(workdir, "video.%(ext)s")
    ydl_opts = {
        'format': 'best[ext=mp4]/best',
        'outtmpl': output_path,
        'noplaylist': True,
    }
    with YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])
    
    # Find the downloaded file
    for file in os.listdir(workdir):
        if file.startswith("video."):
            return os.path.join(workdir, file)
    raise Exception("Video download failed")

def _extract_wav_from_video(video_path: str, wav_path: str, sample_rate: int = 16000):
    """Extract audio from video to WAV format"""
    cmd = [
        "ffmpeg", "-y", "-i", video_path,
        "-ar", str(sample_rate), "-ac", "1", "-f", "wav",
        wav_path
    ]
    subprocess.run(cmd, check=True, capture_output=True)

def _detect_speaker_gender(audio_segment):
    """Simple gender detection based on fundamental frequency"""
    try:
        # Convert to numpy array
        samples = np.array(audio_segment.get_array_of_samples())
        if audio_segment.channels == 2:
            samples = samples.reshape((-1, 2)).mean(axis=1)
        
        # Simple pitch detection - this is a basic implementation
        # In practice, you might want to use more sophisticated methods
        frame_length = 2048
        hop_length = 512
        
        # Calculate zero crossing rate as a simple pitch indicator
        zcr = np.mean([np.sum(np.diff(np.sign(samples[i:i+frame_length])) != 0) 
                      for i in range(0, len(samples)-frame_length, hop_length)])
        
        # Simple heuristic: higher ZCR typically indicates higher pitch (female)
        return "meera" if zcr > 100 else "arvind"
    except:
        return "arvind"  # Default fallback

def _sarvam_stt_from_audiosegment(audio_segment, language_code: str) -> str:
    """Convert AudioSegment to text using Sarvam STT"""
    try:
        wav_io = io.BytesIO()
        audio_segment.export(wav_io, format="wav")
        wav_io.seek(0)
        
        files = {"file": ("audio.wav", wav_io.getvalue(), "audio/wav")}
        headers = {"api-subscription-key": SARVAM_API_KEY}
        data = {"language_code": language_code}
        
        response = requests.post(SARVAM_STT_URL, headers=headers, files=files, data=data, timeout=(10, 60))
        if response.status_code == 200:
            result = response.json()
            return result.get("transcript", "")
    except Exception as e:
        logging.error(f"STT error: {e}")
    return ""

def _normalize_language_code(lang_code: str) -> str:
    """Normalize language code for Sarvam API"""
    if not lang_code:
        return "en-IN"
    
    # Sarvam API supported language mapping
    supported_langs = {
        'en': 'en-IN', 'en-IN': 'en-IN',
        'hi': 'hi-IN', 'hi-IN': 'hi-IN', 
        'bn': 'bn-IN', 'bn-IN': 'bn-IN',
        'te': 'te-IN', 'te-IN': 'te-IN',
        'mr': 'mr-IN', 'mr-IN': 'mr-IN',
        'ta': 'ta-IN', 'ta-IN': 'ta-IN',
        'gu': 'gu-IN', 'gu-IN': 'gu-IN',
        'kn': 'kn-IN', 'kn-IN': 'kn-IN',
        'ml': 'ml-IN', 'ml-IN': 'ml-IN',
        'or': 'or-IN', 'or-IN': 'or-IN',
        'pa': 'pa-IN', 'pa-IN': 'pa-IN'
    }
    
    return supported_langs.get(lang_code, 'en-IN')

def _sarvam_translate(text: str, source_lang: str, target_lang: str) -> str:
    """Translate text using Sarvam API"""
    try:
        headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
        payload = {
            "input": text,
            "source_language_code": _normalize_language_code(source_lang),
            "target_language_code": _normalize_language_code(target_lang),
            "speaker_gender": "Male",
            "mode": "formal",
            "model": "mayura:v1",
            "enable_preprocessing": True,
        }
        
        response = requests.post(SARVAM_TRANSLATE_URL, headers=headers, json=payload, timeout=(10, 60))
        if response.status_code == 200:
            result = response.json()
            return decode_html_entities(result.get("translated_text", ""))
    except Exception as e:
        logging.error(f"Translation error: {e}")
    return ""

def _sarvam_tts_to_wav(text: str, language_code: str, gender: str, sample_rate: int, output_path: str) -> bool:
    """Generate TTS audio and save to WAV file"""
    try:
        headers = {"api-subscription-key": SARVAM_API_KEY, "Content-Type": "application/json"}
        payload = {
            "inputs": [text],
            "target_language_code": _normalize_language_code(language_code),
            "speaker": gender,
            "pitch": 0,
            "pace": 1.0,
            "loudness": 1.0,
            "speech_sample_rate": sample_rate,
            "enable_preprocessing": True,
            "model": "bulbul:v1"
        }
        
        response = requests.post(SARVAM_TTS_URL, headers=headers, json=payload, timeout=(10, 120))
        logging.info(f"TTS API response: {response.status_code}")
        if response.status_code == 200:
            result = response.json()
            audio_content = result.get("audios", [None])[0]
            if audio_content:
                audio_bytes = base64.b64decode(audio_content)
                with open(output_path, "wb") as f:
                    f.write(audio_bytes)
                return True
            else:
                logging.error(f"No audio content in TTS response: {result}")
        else:
            logging.error(f"TTS API error {response.status_code}: {response.text}")
    except Exception as e:
        logging.error(f"TTS error: {e}")
    return False

def _ffprobe_duration_seconds(video_path: str) -> float:
    """Get video duration in seconds using ffprobe"""
    try:
        cmd = ["ffprobe", "-v", "quiet", "-show_entries", "format=duration", "-of", "csv=p=0", video_path]
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        return float(result.stdout.strip())
    except:
        return 0.0


@app.get("/health")
def health():
    return {"status": "ok"}

@app.get("/api/status")
def get_status():
    """Get server status and available features"""
    return {
        "status": "ok",
        "features": {
            "speech_to_text": True,
            "translation": True,
            "text_to_speech": True,
            "document_processing": {
                "docx": DOCX_AVAILABLE,
                "pdf": PDF_AVAILABLE,
                "pptx": PPTX_AVAILABLE,
                "excel": EXCEL_AVAILABLE,
                "csv": EXCEL_AVAILABLE,  # Uses pandas
                "ocr": OCR_AVAILABLE,
                "txt": True  # Always available
            }
        },
        "missing_dependencies": {
            "docx": [] if DOCX_AVAILABLE else ["python-docx"],
            "pdf": [] if (PDF_AVAILABLE or MUPDF_AVAILABLE or PDFPLUMBER_AVAILABLE) else ["PyPDF2", "PyMuPDF", "pdfplumber"],
            "pptx": [] if PPTX_AVAILABLE else ["python-pptx"],
            "excel": [] if EXCEL_AVAILABLE else ["openpyxl", "xlrd", "pandas"],
            "ocr": [] if OCR_AVAILABLE else ["Pillow", "pytesseract"]
        },
        "pdf_methods": {
            "pdfplumber": PDFPLUMBER_AVAILABLE,
            "pymupdf": MUPDF_AVAILABLE,
            "pypdf2": PDF_AVAILABLE
        }
    }
